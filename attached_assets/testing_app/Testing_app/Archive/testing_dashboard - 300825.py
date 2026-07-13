import dash
from dash import dcc, html, Input, Output, State, callback
import pandas as pd
import plotly.graph_objects as go
from dash import dash_table
import plotly.express as px
import gspread
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from dash.dependencies import Input, Output
from pptx.dml.color import RGBColor
from oauth2client.service_account import ServiceAccountCredentials
from plotly import graph_objects as go
import zipfile 
import os
import json



# this is the setup area of the code

# Enable suppressing callback exceptions
# app = dash.Dash(__name__, suppress_callback_exceptions=True)
app = dash.Dash(__name__, external_stylesheets=[
    "https://fonts.googleapis.com/css2?family=Inter&display=swap"
])

# Google Sheets API scope
scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]


# Choose credentials source
if os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON"):
    print("🌐 Using Render env variable")
    creds_dict = json.loads(os.environ["GOOGLE_SERVICE_ACCOUNT_JSON"])
    creds = ServiceAccountCredentials.from_json_keyfile_dict(creds_dict, scope)
elif os.path.exists("service_account.json"):  # 👈 updated file name here
    print("🖥️ Using local service_account.json")
    creds = ServiceAccountCredentials.from_json_keyfile_name("service_account.json", scope)
else:
    raise EnvironmentError("❌ No credentials found. Set env var or add service_account.json.")

# Authorize client
client = gspread.authorize(creds)


# Load Testing Data
spreadsheet_1sts = client.open("2025-NPLW-1sts-season")
testing_data_1sts = pd.DataFrame(spreadsheet_1sts.worksheet("testing-percentiles1").get_all_records())

button_style = {
    "backgroundColor": "skyblue",
    "color": "black",
    "border": "none",
    "padding": "10px 16px",
    "marginRight": "10px",
    "borderRadius": "6px",
    "fontWeight": "bold",
    "fontFamily": "Arial",
    "cursor": "pointer",
    "boxShadow": "2px 2px 5px rgba(0, 0, 0, 0.3)",
    "textAlign": "left"
}

# Dash Layout

app.layout = html.Div([
    html.Img(src="/assets/clublogo.png", style={"height": "100px", "float": "right", "marginTop": "10px"}),
    html.H1("NPLW - Athletic Testing Dashboard", style={"textAlign": "center", "color": "#FFFFFF", "fontFamily": "Arial Black"}),
    html.Br(),
    html.Br(),

    # Player Report Section
        html.Div([
            html.Label("", style={"color": "#FFFFFF", "fontSize": "14px", "fontFamily": "Arial Black"}),
            dcc.Dropdown(
                id="athletic-player-selector",
                options=[{"label": name, "value": name} for name in sorted(testing_data_1sts["Player Name"].unique())],
                placeholder="Choose a player",
                style={"width": "40%", "margin": "auto", "marginBottom": "10px", "fontFamily": "Arial Black"}
            ),
            html.Div(
            id="athletic-player-report",
            style={
                "color": "white",
                "fontFamily": "Arial Black",
                "padding": "20px",
                "backgroundColor": "#1E3A5F",
                "margin": "auto",
                "width": "80%",
                "whiteSpace": "pre-wrap",
                "border": "1px solid white",
                "borderRadius": "10px"
            }
            ),

        html.Div([
        html.Button("Export Report", id="export-player-report", n_clicks=0, style=button_style),
    dcc.Download(id="download-player-report")
    ], style={"textAlign": "center", "marginTop": "20px"}),    
   
        html.Div([
        html.Button("Export All Reports", id="export-all-reports", n_clicks=0, style=button_style),
        
        # Confirmation popup
        dcc.ConfirmDialog(
            id="confirm-export-all",
            message="Are you sure you want to export reports for all players?"
        ),

        # The download component
        dcc.Download(id="download-all-reports")
    ], style={"textAlign": "center", "marginTop": "10px"}),

    
    #    html.Div([
    #    html.Button("Export All Reports", id="export-all-reports", n_clicks=0, style=button_style),
    #    dcc.Download(id="download-all-reports")
    #], style={"textAlign": "center", "marginTop": "10px"}),



    html.Br(),
    html.Br(),
    html.Br(),

    # Athletic Testing Overview
    html.Div([
        html.H2("Athletic Testing Overview", style={"textAlign": "center", "color": "white", "fontFamily": "Arial Black"}),

        html.Div([
            html.Label("", style={"color": "#FFFFFF", "fontSize": "18px"}),
            dcc.Dropdown(
                id="athletic-metric-selector",
                options=[
                    {"label": "Vertical Total", "value": "Vertical total"},
                    {"label": "Vertical Relative", "value": "Vertical relative"},
                    {"label": "Horizontal", "value": "Horizontal"},
                    {"label": "Balsom", "value": "Balsom"},
                    {"label": "0–10m Split", "value": "0-10 split"},
                    {"label": "10–20m Split", "value": "10-20 split"},
                    {"label": "20–30m Split", "value": "20-30 split"},
                ],
                value="Vertical total",
                clearable=False,
                #style={"width": "40%", "margin": "auto", "backgroundColor": "#FFFFFF", "color": "#000000", "fontFamily": "Arial Black"}
                style={"width": "40%", "margin": "auto", "marginBottom": "10px", "fontFamily": "Arial Black"}
            )
        ], style={"textAlign": "center", "padding": "10px"}),

        html.Div([
        html.Button("Sort High to Low", id="sort-high-testing", n_clicks=0, style=button_style),
        html.Button("Sort Low to High", id="sort-low-testing", n_clicks=0, style=button_style)
        
    ], style={"textAlign": "left", "padding": "10px"}),

        dcc.Graph(id="athletic-testing-chart", style={"backgroundColor": "black"}),
    #was #ccc in solid white
    ], style={"backgroundColor": "#1E3A5F", "padding": "20px", "border": "1px solid white", "borderRadius": "10px"}),

    html.Br(),
    html.Br(),
    html.Br(),

    # Second testing chart tier grouping layout
    html.Div([
        html.Div([
            html.Label("", style={"color": "#FFFFFF", "fontSize": "14px", "fontFamily": "Arial Black"}),
            dcc.Dropdown(
                id="tier-metric-selector",
                options=[
                    {"label": "Vertical total", "value": "Vertical total"},
                    {"label": "Vertical relative", "value": "Vertical relative"},
                    {"label": "Horizontal", "value": "Horizontal"},
                    {"label": "Balsom", "value": "Balsom"},
                    {"label": "0-10 split", "value": "0-10 split"},
                    {"label": "10-20 split", "value": "10-20 split"},
                    {"label": "20-30 split", "value": "20-30 split"},
                ],
                value="Vertical total",
                clearable=False,
                style={"width": "35%", "margin": "auto", "marginBottom": "10px", "fontFamily": "Arial Black"}
            )
        ], style={"textAlign": "center", "paddingBottom": "10px"}),

        html.Label("This chart will show highest to lowest for each testing element.", style={
            "color": "white",
            "fontWeight": "bold",
            "textAlign": "left",
            "fontFamily": "Arial",
            "marginBottom": "5px"
        }),
        
        dcc.Graph(id="tier-grouping-chart", style={"backgroundColor": "black", "marginTop": "15px"}),
    ], style={"backgroundColor": "#1E3A5F", "padding": "20px", "border": "1px solid white", "borderRadius": "10px"}),

    html.Br(),
    html.Br(),

    # filtering by tier line chart
    html.Div([
    html.Div([
    
    html.Label("", style={"color": "#FFFFFF", "fontSize": "14px", "fontFamily": "Arial Black"}),
    dcc.Dropdown(
        id="sprint-tier-filter",
        options=[
            {"label": "All Tiers", "value": "All"},
            {"label": "Green", "value": "Green"},
            {"label": "Light Green", "value": "Light Green"},
            {"label": "Orange", "value": "Orange"},
            {"label": "Red", "value": "Red"}
        ],
        value="All",
        clearable=False,
        style={
            "width": "35%",
            "margin": "auto",
            "marginBottom": "10px",
            "fontFamily": "Arial Black"
        }
    )
], style={"textAlign": "center", "paddingBottom": "10px"}),

    html.Label("Filter by Tier (based on 0–10 split). You can see from their first split if they stay steady, accelerate or decline through over the distance.\n", style={
        "color": "white",
        "fontWeight": "bold",
        "fontFamily": "Arial",
        "marginBottom": "5px"
    }),
    

    dcc.Graph(id="sprint-line-chart", style={"backgroundColor": "black", "marginTop": "15px"}),
    ], style={"backgroundColor": "#1E3A5F", "padding": "20px", "border": "1px solid white", "borderRadius": "10px"}),
      
    ]),
], style={"backgroundColor": "#1E3A5F", "padding": "20px", "minHeight": "100vh"})  # ✅ BLUE BACKGROUND FIX
      # this will always be the last line of the layouts area, it closes the fist [ and then defines the parameters of `app.layout = html.Div([...])`


# callbacks start here.

# function created to be reused by more than one callback
def build_athletic_report(player_row, player_name):
    def parse(val):
        try:
            return float(val)
        except:
            return 0.0



    def std(val):
        return 100 - val

    vertical_total = parse(player_row["Vertical total"].values[0])
    vertical_relative = parse(player_row["Vertical relative"].values[0])
    horizontal = parse(player_row["Horizontal"].values[0])
    balsom = std(parse(player_row["Balsom"].values[0]))
    split_0_10 = std(parse(player_row["0-10 split"].values[0]))
    split_10_20 = std(parse(player_row["10-20 split"].values[0]))
    split_20_30 = std(parse(player_row["20-30 split"].values[0]))

    # Same band, comment, and sprint logic goes here…
    # (just copy/paste your current inner functions)

    def band(val):
        if val == 0:
            return "No data"
        elif val >= 66:
            return "🟢 "
        elif val <= 33:
            return "🔴 "
        else:
            return "🟠 "

    def jump_comment(val):
        if val == 0:
            return "No data to assess."
        advice = []
        if val >= 50:
            advice.append("\nYou should expect to be involved in attacking and defensive set pieces close to goal.")
            if val >= 66:
                advice.append("\nTactical – You’re strong in the air. Be brave, attack the ball early, and take up aggressive positions near goal during set pieces.")
        else:
            advice.append("\nTraining – Build your jump power with exercises like box jumps and squat jumps. This will help in aerial duels.")
            if val <= 33:
                advice.append("\nTactical – If you're not the biggest jumper, use smart positioning and time your jump well. You can still win the contest.")
        return " ".join(advice)


    def vertical_explosiveness_comment(val):
        if val == 0:
            return "No data to assess."
        advice = []
        if val >= 66:
            advice.append("\nTraining – Your vertical explosiveness is strong. Keep it sharp with reactive jumps, single-leg plyos, and low-rep strength work.")
            advice.append("\nTactical – You’ve got a strong take-off. Use it to win space early in 1v1 contests and pressing moments.")
        elif val <= 33:
            advice.append("\nTraining – Work on power with contrast training: squats plus jumps, resistance jumps, and med ball throws.")
            advice.append("\nTactical – If you're behind on power, make up for it by reading the play early and getting into good spots.")
        else:
            advice.append("\nTraining – You're at a good level. Maintain it with pogo hops, snap jumps, or short jump series.")
            advice.append("\nTactical – You can hold your own in contests. Back yourself in duels with opponents.")
        return " ".join(advice)


    def explosiveness_comment(val):
        if val == 0:
            return "No data to assess."
        advice = []
        if val <= 33:
            advice.append("\nTraining – Improve your explosiveness with bounding drills, resistance sprints, and quick box jumps.")
            advice.append("\nTactical – If others are quicker, make the first move. Good positioning gives you a head start.")
        elif val >= 66:
            advice.append("\nTraining – Your explosiveness is strong. Keep it sharp with reactive drills and occasional power work.")
            advice.append("\nTactical – Use your power to burst past players in attack or close down quickly in defence.")
        else:
            advice.append("\nTraining – You’re solid here. Maintain with reactive power work like med ball throws and plyo bounds.")
            advice.append("\nTactical – You’ve got a reliable burst. Use it in tight spaces to press, pass, or drive forward.")
        return " ".join(advice)

        

    def agility_comment(val):
        if val == 0:
            return "No data to assess."
        advice = []
        if val >= 66:
            advice.append("\nTraining – You’ve got strong agility — no extra work needed right now.")
            advice.append("\nTactical – You're capable of defending and attacking in tight spaces. Trust your agility to stay tight in duels and escape pressure confidently.")
        elif val <= 33:
            advice.append("\nTraining - You can improve your agility and change of direction with drills like cone weaving, 5-10-5 runs, and shuttle sprints. Focus on staying low and controlled through the turns.")
            advice.append("\nTactical - In tight defensive moments, smart positioning and quick decision-making can help you compete with players who are more agile. Focus on anticipation and body orientation.")
        else:
            advice.append("\nTraining - Your agility is solid. Maintain it with quick-footed drills like ladder work, short shuttle runs, and small-space directional changes.")
            advice.append("\nTactical – Your agility gives you a solid base in tight areas. You can stay competitive in small spaces, especially when you anticipate early and stay light on your feet.")
        return " ".join(advice)

    def band_sprint(val):
        if val == 0:
            return "⚪"
        elif val >= 66:
            return "🟢"
        elif val <= 33:
            return "🔴"
        else:
            return "🟠"

    def sprint_comment(s0, s10, s20):
        if s0 == 0 or s10 == 0 or s20 == 0:
            return "No data to assess."
        notes = []
        if s0 <= 35:
            notes.append("\nSplit 1 (0-10m) comments: \nTraining – Improve your first-step acceleration with sled pushes, incline sprints, and resisted acceleration drills.")
            notes.append("\nTactical - Be mindful in 1v1 moments where quick reactions are needed — focus on good positioning and make decisions early to compensate.")
        elif s0 >= 66:
            notes.append("\nSplit 1 (0-10m) comments: \nTraining – Your acceleration is strong. You don’t need targeted work here — just maintain with sharp technical sprint starts.")
            notes.append("\nTactical - In possession you can be confident that you can burst past opponents. In defence, remember you can close players down quickly, or react to their changes in direction.")

        if s20 <= 33:
            notes.append("\nSplit 3 (20-30m) comments: \nTraining - You can make improvements to max velocity, by doing exercises like flying sprints, speed endurance.")
            notes.append("\nTactical - For longer sprint moments, try to read the play early. Intelligent positioning can help you avoid footraces you may not win.")
        elif s20 >= 66:
            notes.append("\nSplit 3 (20-30m) comments: \nTraining - No specific speed-endurance focus required, you are strong in this area.")
            notes.append("\nTactical – You're well suited to chasing through balls or covering large spaces defensively — use your top-end speed to your advantage in transition moments.")
        
        if not notes:
            notes.append("\nGeneral comments: \nTraining - Your sprint profile is balanced, placing you around the team average. Maintain or improve this with technical sprint drills like flying sprints, sprint-float-sprint, and resisted band starts.")
            notes.append("\nGeneral comments: \nTactical – You can compete in both short and long sprint moments. Use this balance to support play on both sides of the ball, especially in transition moments.")
        return " ".join(notes)

# **{player_name}'s Testing Report – From 20 Mar 2025 Testing** this was under report = f"""
    report = f"""

🦘 Jumping Ability
 • Jump Height: {round(vertical_total)}% – {band(vertical_total)} {jump_comment(vertical_total)}
\n• Explosive Power: {round(vertical_relative)}% – {band(vertical_relative)} {vertical_explosiveness_comment(vertical_relative)}
\n• Bounding Power: {round(horizontal)}% – {band(horizontal)} {explosiveness_comment(horizontal)}

🌀 Agility
  • Agility Score: {round(balsom)}% – {band(balsom)} {agility_comment(balsom)}

⚡ Sprint Profile
  • 0–10m: {round(split_0_10)}% {band_sprint(split_0_10)}, 10–20m: {round(split_10_20)}% {band_sprint(split_10_20)}, 20–30m: {round(split_20_30)}% {band_sprint(split_20_30)} {sprint_comment(split_0_10, split_10_20, split_20_30)}

""".strip()


    return report

# are you sure to create all reports? callback
@callback(
    Output("confirm-export-all", "displayed"),
    Input("export-all-reports", "n_clicks"),
    prevent_initial_call=True
)
def show_confirmation(n_clicks):
    return True  # show popup


#Athletic testing overview
@app.callback(
    Output("athletic-testing-chart", "figure"),
    [Input("athletic-metric-selector", "value"),
     Input("sort-high-testing", "n_clicks"),
     Input("sort-low-testing", "n_clicks")]
)
def update_athletic_testing_chart(selected_metric, high_clicks, low_clicks):
    testing_data = testing_data_1sts.copy()

    print("Available columns:", testing_data.columns.tolist())
    print("Selected metric:", selected_metric)
    print(testing_data[[selected_metric, "Player Name"]].head())

    # Handle missing or non-numeric entries
    testing_data[selected_metric] = pd.to_numeric(testing_data[selected_metric], errors="coerce")

    # Define metrics that are reversed (lower percentile = better)
    reverse_metrics = ["Balsom", "0-10 split", "10-20 split", "20-30 split"]

    def get_color(val, metric):
        try:
            val = float(val)
            if val == 0:
                return "gray"  # Don't color if value is 0
            if metric in ["Balsom", "0-10 split", "10-20 split", "20-30 split"]:
                # Lower is better
                if val >= 66:
                    return "red"
                elif val >= 40:
                    return "orange"
                elif val >= 20:
                    return "lightgreen"
                else:
                    return "green"
            else:
                # Higher is better
                if val >= 80:
                    return "green"
                elif val >= 60:
                    return "lightgreen"
                elif val >= 40:
                    return "orange"
                else:
                    return "red"
        except:
            return "gray"

    def get_training_note(val, selected_metric):
        try:
            val = float(val)
            if val == 0:
                return "N/A"
            if selected_metric in ["Balsom", "0-10 split", "10-20 split", "20-30 split"]:
                if val >= 66:
                    return "Needs improvement"
                elif val >= 40:
                    return "Solid"
                else:
                    return "Strong"
            else:
                if val <= 40:
                    return "Needs improvement"
                elif val <= 60:
                    return "Solid"
                else:
                    return "Strong"
        except:
            return "Unknown"


    def get_training_focus(val, selected_metric):
        try:
            val = float(val)
            if val == 0:
                return "No data to assess"
            if selected_metric == "Vertical total":
                if val >= 50:
                    return "Should be in the box for set pieces"
                else:
                    return "Not expected to win aerial duels"
            elif selected_metric in ["Vertical relative", "Horizontal"]:
                if val <= 33:
                    return "Focus on improving explosiveness"
                elif val >= 66:
                    return "No specific focus needed"
                else:
                    return "Solid explosiveness"
            elif selected_metric == "Balsom":
                if val >= 66:
                    return "Work on agility and change of direction"
                elif val <= 33:
                    return "No specific focus needed"
                else:
                    return "Solid agility"
            elif selected_metric in ["0-10 split", "10-20 split", "20-30 split"]:
                return "Check split pattern across all phases"
            else:
                return "General improvement area"
        except:
            return "Unknown"



    # Add Performance Tier text
    testing_data[selected_metric] = pd.to_numeric(testing_data[selected_metric], errors="coerce").fillna(0)
    #testing_data["Color"] = testing_data[selected_metric].apply(lambda val: get_color(val, selected_metric))
    #testing_data["Tier"] = testing_data[selected_metric].apply(lambda val: get_training_note(val, selected_metric))
    #testing_data["Training Focus"] = testing_data[selected_metric].apply(lambda val: get_training_focus(val, selected_metric))
    testing_data["Color"] = testing_data.apply(lambda row: get_color(row[selected_metric], selected_metric), axis=1)
    testing_data["Tier"] = testing_data.apply(lambda row: get_training_note(row[selected_metric], selected_metric), axis=1)
    testing_data["Training Focus"] = testing_data.apply(lambda row: get_training_focus(row[selected_metric], selected_metric), axis=1)


    print(testing_data[["Player Name", selected_metric, "Tier", "Training Focus"]])



    # Sort
    if high_clicks > low_clicks:
        testing_data = testing_data.sort_values(by=selected_metric, ascending=False)
    elif low_clicks > high_clicks:
        testing_data = testing_data.sort_values(by=selected_metric, ascending=True)

    fig = px.bar(
        testing_data,
        x="Player Name",
        y=selected_metric,
        color="Color",
        color_discrete_map={
            "red": "red",
            "orange": "orange",
            "lightgreen": "lightgreen",
            "green": "green",
            "gray": "gray"
        },
        title=f"Athletic Testing – {selected_metric}",
        template="plotly_dark"
    )

    fig.update_traces(
        hovertemplate="<b>%{x}</b><br>Percentile: %{y}%<br>Performance: %{customdata[0]}<br>Training Focus: %{customdata[1]}",
        customdata=testing_data[["Tier", "Training Focus"]]
    )

    return fig

# Performance tiers chart
@app.callback(
    Output("tier-grouping-chart", "figure"),
    [Input("tier-metric-selector", "value")]
)
def update_tier_grouping_chart(selected_metric):
    df = testing_data_1sts.copy()
    print("Selected Metric:", selected_metric)
    print("Available Columns:", df.columns.tolist())


    # Convert to numeric
    df[selected_metric] = pd.to_numeric(df[selected_metric], errors="coerce").fillna(0)
    print(df[["Player Name", selected_metric]].head(10))


    # Define reverse-logic metrics (lower = better)
    reverse_metrics = ["Balsom", "0-10 split", "10-20 split", "20-30 split"]

    # Tier classification
    def get_tier(val):
        try:
            val = float(val)
            if val == 0:
                return "No Data"
            if selected_metric in reverse_metrics:
                if val >= 66:
                    return "Red"
                elif val >= 40:
                    return "Light Red"
                elif val >= 20:
                    return "Light Green"
                else:
                    return "Green"
            else:
                if val >= 80:
                    return "Green"
                elif val >= 60:
                    return "Light Green"
                elif val >= 40:
                    return "Light Red"
                else:
                    return "Red"
        except:
            return "No Data"

    df["Tier"] = df[selected_metric].apply(get_tier)
    print(df[["Player Name", selected_metric, "Tier"]].head(10))


    # Sorting for display
    ascending = selected_metric in reverse_metrics
    df = df.sort_values(by=selected_metric, ascending=ascending)

    # Define tier order
    tier_order = ["Green", "Light Green", "Light Red", "Red", "No Data"]

    fig = px.bar(
        df,
        x="Player Name",
        y=selected_metric,
        color="Tier",
        color_discrete_map={
            "Green": "green",
            "Light Green": "lightgreen",
            "Light Red": "orange",
            "Red": "red",
            "No Data": "gray"
        },
        category_orders={"Tier": tier_order},
        title=f"Player Performance Tiers – {selected_metric}",
        template="plotly_dark"
    )

    fig.update_traces(
        hovertemplate="<b>%{x}</b><br>Percentile: %{y}%<br>Tier: %{customdata}",
        customdata=df["Tier"]
    )

    return fig

#player report paragraph
@callback(
    Output("athletic-player-report", "children"),
    Input("athletic-player-selector", "value")
)
def generate_player_report(player_name):
    if not player_name:
        return ""

    df = testing_data_1sts.copy()
    df.replace("", 0, inplace=True)
    player_row = df[df["Player Name"] == player_name]
    
    if player_row.empty:
        return f"No data found for {player_name}."

    report = build_athletic_report(player_row, player_name)
    # return dcc.Markdown(report)
    return dcc.Markdown(f"**{player_name}'s Testing Report – From 20 Mar 2025 Testing**\n\n{build_athletic_report(player_row, player_name)}")



#callback for report download
@callback(
    Output("download-player-report", "data"),
    Input("export-player-report", "n_clicks"),
    State("athletic-player-selector", "value"),
    prevent_initial_call=True
)
def export_player_report(n_clicks, player_name):
    if not player_name:
        return None

    df = testing_data_1sts.copy()
    df.replace("", 0, inplace=True)
    row = df[df["Player Name"] == player_name]
    if row.empty:
        return None

    # 📝 Generate full paragraph from shared function
    report_text = build_athletic_report(row, player_name)

    # 🖼️ Create PowerPoint
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]

    # 🎨 Style setup
    dark_blue = RGBColor(15, 31, 51)
    light_blue = RGBColor(194, 224, 241)

    # ✂️ Split report into two parts
    section_split = "🌀"
    if section_split in report_text:
        jumping_text, agility_sprint_text = report_text.split(section_split, 1)
        agility_sprint_text = f"{section_split}{agility_sprint_text.strip()}"
    else:
        jumping_text = report_text
        agility_sprint_text = ""

    def add_slide(title_text, content_text, include_logo=False):
        slide = prs.slides.add_slide(slide_layout)

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = light_blue

        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_tf = title.text_frame
        title_tf.text = title_text
        title_tf.paragraphs[0].font.size = Pt(24)
        title_tf.paragraphs[0].font.name = "Segoe UI Black"
        title_tf.paragraphs[0].font.color.rgb = dark_blue
        title_tf.paragraphs[0].alignment = 1  # center

        # Logo (optional)
        if include_logo:
            try:
                slide.shapes.add_picture("assets/clublogo.png", Inches(7.5), Inches(0.2), width=Inches(1.3))
            except Exception as e:
                print("Logo not found or failed to load:", e)

        # Content
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.2), Inches(6))
        tf = box.text_frame
        tf.word_wrap = True

        for line in content_text.strip().splitlines():
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.name = "Segoe UI"
            p.font.color.rgb = dark_blue

    # Slide 1 – Jumping
    add_slide(f"{player_name} – Testing Report – March 2025", jumping_text, include_logo=True)

    # Slide 2 – Agility + Sprint (only if exists)
    if agility_sprint_text.strip():
        add_slide(f"{player_name} – Testing Report – March 2025", agility_sprint_text, include_logo=True)

    # Export
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return dcc.send_bytes(pptx_io.read(), f"{player_name}_Athletic_Report.pptx")



# no idea what this callback is for
@callback(
    Output("sprint-line-chart", "figure"),
    Input("sprint-tier-filter", "value")
)
def update_sprint_line_chart(selected_tier):
    df = testing_data_1sts.copy()
    df.replace("", 0, inplace=True)

    # Step 1: Convert to numeric, blanks become NaN
    for col in ["0-10 split", "10-20 split", "20-30 split"]:
        df[col] = pd.to_numeric(df[col], errors="coerce")

    # Step 2: Drop rows with any NaN in the split columns BEFORE inverting
    df = df.dropna(subset=["0-10 split", "10-20 split", "20-30 split"]).copy()

    # Step 3: Invert percentiles so higher = better (true inversion)
    for col in ["0-10 split", "10-20 split", "20-30 split"]:
        df[col] = 100 - df[col]


    # Now filter out rows with no valid data
    df = df[
        (df["0-10 split"] > 0) &
        (df["10-20 split"] > 0) &
        (df["20-30 split"] > 0)
    ].copy()

    # Tier classification based on 0–10 split
    def classify_tier(val):
        if val >= 66:
            return "Green"
        elif val >= 40:
            return "Light Green"
        elif val >= 20:
            return "Orange"
        else:
            return "Red"

    df["Tier"] = df["0-10 split"].apply(classify_tier)

    # Filter by dropdown selection
    if selected_tier != "All":
        df = df[df["Tier"] == selected_tier].copy()


    # Debug prints (optional)
    print("Selected tier:", selected_tier)
    print("Tiers present:", df["Tier"].unique())
    print("Row count after filtering:", len(df))


    # Create figure
    colors = px.colors.qualitative.Light24
    fig = go.Figure(data=[])  # Clear out any previous traces


    for i, (_, row) in enumerate(df.iterrows()):
        fig.add_trace(go.Scatter(
            x=["0–10", "10–20", "20–30"],
            y=[round(row["0-10 split"]), round(row["10-20 split"]), round(row["20-30 split"])],
            mode="lines+markers",
            name=row["Player Name"],
            line=dict(width=2, color=colors[i % len(colors)]),
            hovertemplate=(
                f"<b>{row['Player Name']}</b><br>"
                f"0–10 split: {round(row['0-10 split'])}%<br>"
                f"10–20 split: {round(row['10-20 split'])}%<br>"
                f"20–30 split: {round(row['20-30 split'])}%"
            )

           
        ))

    fig.update_layout(
        title="Sprint Split Profile by Player",
        xaxis_title="Split Phase",
        yaxis_title="Percentile (Inverted – higher is better)",
        plot_bgcolor="black",
        paper_bgcolor="black",
        font=dict(color="white", family="Arial"),
        hovermode="closest",
        xaxis=dict(gridcolor="gray"),  # This colours gridlines grey
        yaxis=dict(showgrid=False)  # removes horizontal lines
    )


    return fig


# export all reports callback for export all reports button
@callback(
    Output("download-all-reports", "data"),
    Input("confirm-export-all", "submit_n_clicks"),
    prevent_initial_call=True
)
def export_all_reports(n_clicks):
    if n_clicks == 0:
        return None

    df = testing_data_1sts.copy()
    df.replace("", 0, inplace=True)

    # PowerPoint setup
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]

    dark_blue = RGBColor(15, 31, 51)
    light_blue = RGBColor(194, 224, 241)

    def add_slide(title_text, content_text, include_logo=False):
        slide = prs.slides.add_slide(slide_layout)

        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = light_blue

        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_tf = title.text_frame
        title_tf.text = title_text
        title_tf.paragraphs[0].font.size = Pt(24)
        title_tf.paragraphs[0].font.name = "Segoe UI Black"
        title_tf.paragraphs[0].font.color.rgb = dark_blue
        title_tf.paragraphs[0].alignment = 1  # Centered

        # Optional logo
        if include_logo:
            try:
                slide.shapes.add_picture("assets/clublogo.png", Inches(7.5), Inches(0.2), width=Inches(1.3))
            except Exception as e:
                print("Logo not found or failed to load:", e)

        # Content box
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.2), Inches(6))
        tf = box.text_frame
        tf.word_wrap = True

        for line in content_text.strip().splitlines():
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.name = "Segoe UI"
            p.font.color.rgb = dark_blue

    # Loop through players
    for player_name in df["Player Name"].unique():
        player_row = df[df["Player Name"] == player_name]
        if player_row.empty:
            continue

        report_text = build_athletic_report(player_row, player_name)

        # Split into two sections
        section_split = "🌀"
        if section_split in report_text:
            jumping_text, agility_sprint_text = report_text.split(section_split, 1)
            agility_sprint_text = f"{section_split}{agility_sprint_text.strip()}"
        else:
            jumping_text = report_text
            agility_sprint_text = ""

        # Add slide 1 – Jumping
        add_slide(f"{player_name} – Testing Report – March 2025", jumping_text, include_logo=True)

        # Add slide 2 – Agility + Sprint
        if agility_sprint_text.strip():
            add_slide(f"{player_name} – Testing Report – March 2025", agility_sprint_text, include_logo=True)

    import zipfile

    # Create ZIP buffer
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "w") as zip_file:
        for _, player_row in testing_data_1sts.iterrows():
            player_name = player_row["Player Name"]
            report_text = build_athletic_report(player_row.to_frame().T, player_name)

            # Split report content
            section_split = "🌀"
            if section_split in report_text:
                jumping_text, agility_sprint_text = report_text.split(section_split, 1)
                agility_sprint_text = f"{section_split}{agility_sprint_text.strip()}"
            else:
                jumping_text = report_text
                agility_sprint_text = ""

            # Create PPTX for this player
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]

            def add_slide(title_text, content_text, include_logo=False):
                slide = prs.slides.add_slide(slide_layout)
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(194, 224, 241)

                title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                tf = title.text_frame
                tf.text = title_text
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.name = "Segoe UI Black"
                tf.paragraphs[0].font.color.rgb = RGBColor(15, 31, 51)
                tf.paragraphs[0].alignment = 1

                if include_logo:
                    try:
                        slide.shapes.add_picture("assets/clublogo.png", Inches(7.5), Inches(0.2), width=Inches(1.3))
                    except Exception as e:
                        print(f"Logo error for {player_name}:", e)

                box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.2), Inches(6))
                tf = box.text_frame
                tf.word_wrap = True
                for line in content_text.strip().splitlines():
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(16)
                    p.font.name = "Segoe UI"
                    p.font.color.rgb = RGBColor(15, 31, 51)

            add_slide(f"{player_name} – Testing Report – March 2025", jumping_text, include_logo=True)
            if agility_sprint_text.strip():
                add_slide(f"{player_name} – Testing Report – March 2025", agility_sprint_text, include_logo=True)

            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)

            zip_file.writestr(f"{player_name}_Athletic_Report_25-03-2025_Testing.pptx", pptx_bytes.read())

    zip_buffer.seek(0)
    return dcc.send_bytes(zip_buffer.read(), "All_Player_Reports.zip")

# Run the app
# if __name__ == '__main__':
    # Use this line when running on your home network
    # app.run_server(debug=True, host='192.168.20.21', port=8050)
    
    # Use this line when off the home network
    # app.run_server(debug=True, port=8050)

    # use this if running from render
    # app.run(host='0.0.0.0', port=8050, debug=False)


    #app.run_server(debug=True)

    #this is the running version on the web

if __name__ == "__main__":
    app.run(
        host="0.0.0.0",
        port=int(os.environ.get("PORT", 8050)),
        debug=True
    )